"""生成项目介绍 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x09, 0x09, 0x0B)
BG_CARD = RGBColor(0x12, 0x12, 0x1A)
ACCENT = RGBColor(0x81, 0x8C, 0xF8)
ACCENT2 = RGBColor(0x63, 0x66, 0xF1)
TEXT = RGBColor(0xD4, 0xD4, 0xD8)
TEXT_DIM = RGBColor(0x71, 0x71, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide

def add_title(slide, text, left=0.5, top=0.3, width=12, size=36):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    return tf

def add_subtitle(slide, text, left=0.5, top=1.1, width=12, size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = TEXT_DIM
    return tf

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def card_text(shape, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet(tf, text, level=0, size=13, color=TEXT_DIM):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return p

# ====== SLIDE 1: 封面 ======
s = dark_slide()
add_title(s, "Python 多功能智能助手", top=2.0, size=48)
add_subtitle(s, "基于 Flask + LangChain 的多模型 AI 对话系统", top=2.8, size=22)
add_subtitle(s, "项目技术汇报", top=3.4, size=16)

# ====== SLIDE 2: 项目概述 ======
s = dark_slide()
add_title(s, "项目概述", size=32)
add_subtitle(s, "一个支持语音对话、多角色切换、主题自定义、文件上传的智能 AI 聊天助手", size=16)

items = [
    ("多模型支持", "DeepSeek + 通义千问 双模型，统一接口动态切换"),
    ("流式响应", "SSE Token 级实时输出，Markdown 边出边排版"),
    ("语音对话", "浏览器 STT 语音输入 + 火山/腾讯双引擎 TTS，11 种音色"),
    ("角色系统", "6 内置角色 + 无限自定义，对话-角色绑定，历史自动恢复"),
    ("主题系统", "12 套预设主题，全透明度调节，背景图/视频自定义"),
    ("附件上传", "支持 PDF/Word/TXT/代码文件，AI 基于内容回答"),
]
for i, (title, desc) in enumerate(items):
    row, col = divmod(i, 3)
    c = add_card(s, 0.5+col*4.1, 2.0+row*2.5, 3.8, 2.0)
    tf = card_text(c, title, size=16, color=ACCENT, bold=True)
    add_bullet(tf, desc, size=12)

# ====== SLIDE 3: 技术栈 ======
s = dark_slide()
add_title(s, "技术栈总览", size=32)
add_subtitle(s, "前后端 + 数据库 + AI + 部署 全链路技术选型", size=16)

techs = [
    ("后端框架", "Flask 3.x, Flask-SocketIO, Flask-CORS"),
    ("AI/LLM", "LangChain 1.x, DeepSeek API, 通义千问 DashScope"),
    ("数据库", "MySQL 8.0 + SQLAlchemy ORM, Redis 7.0 (限流)"),
    ("语音", "Web Speech API (STT), 火山引擎 TTS, 腾讯云 TTS"),
    ("前端", "原生 HTML/CSS/JavaScript, SSE, marked.js, highlight.js"),
    ("部署", "Gunicorn + gevent + Nginx, Docker Compose"),
]
for i, (title, items) in enumerate(techs):
    row, col = divmod(i, 3)
    c = add_card(s, 0.5+col*4.1, 2.0+row*2.5, 3.8, 1.8)
    tf = card_text(c, title, size=16, color=ACCENT2, bold=True)
    for item in items.split(", "):
        add_bullet(tf, "· "+item, size=12)

# ====== SLIDE 4: 系统架构 ======
s = dark_slide()
add_title(s, "系统架构", size=32)
add_subtitle(s, "分层架构设计：API路由层 → 业务服务层 → LLM调用层 → 数据持久层", size=16)

arch_text = """
浏览器 (语音STT + 粒子UI + SSE流式)
    │  HTTP / SSE / WebSocket
    ▼
Nginx (生产) → Gunicorn + gevent
    │
    ▼
Flask 应用工厂
    ├── API 路由层 (5 蓝图, 20+ 接口)
    │     chat / conversation / role / voice / system / upload
    ├── 业务服务层
    │     ChatService / ModelService / ToolService
    ├── LLM 调用层 (策略模式)
    │     DeepSeekClient / QwenClient / Memory
    ├── 中间件层
    │     Auth / RateLimit (Redis) / Logger
    └── 数据模型层
          Conversation / Message / Role / User
              │
              ▼
          MySQL 8.0 + Redis 7.0
"""
c = add_card(s, 0.5, 2.0, 12.3, 5.0)
tf = card_text(c, arch_text.strip(), size=12, color=TEXT)

# ====== SLIDE 5: 目录结构 ======
s = dark_slide()
add_title(s, "项目目录结构", size=32)
add_subtitle(s, "40+ Python 文件，模块化分层设计", size=16)

dir_text = """PythonAssistant/
├── run.py                       # 启动入口
├── requirements.txt             # 依赖清单
├── .env                         # 环境变量 (API Key)
├── gunicorn_config.py           # 生产 WSGI 配置
├── Dockerfile / docker-compose.yml
├── app/
│   ├── __init__.py              # 应用工厂 + 角色初始化
│   ├── extensions.py            # DB + SocketIO 扩展
│   ├── config/                  # settings + prompts
│   ├── models/                  # 4 张数据库表
│   │     conversation / message / role / user
│   ├── llm/                     # DeepSeek + Qwen 客户端
│   ├── services/                # chat / model / tool 服务
│   ├── api/                     # 6 个蓝图 (20+ API)
│   ├── tools/                   # 计算器/数据分析/搜索
│   ├── middleware/              # 认证/限流/日志
│   └── utils/                   # 响应/异常/校验
├── static/                      # 前端 SPA + 简历 + 麦克风测试
├── tests/                       # 单元测试
└── docs/                        # README / deploy / 零基础指南"""
c = add_card(s, 0.5, 2.0, 12.3, 5.0)
tf = card_text(c, dir_text, size=11, color=TEXT)

# ====== SLIDE 6: 核心功能 - 对话流程 ======
s = dark_slide()
add_title(s, "核心功能：对话流程", size=32)
add_subtitle(s, "从用户输入到 AI 流式回复的完整链路", size=16)

flow_text = """1. 用户打字/语音 → 前端 fetch POST /api/v1/chat/stream
2. Flask 接收请求 → ChatService.prepare_chat_context()
   - 加载角色 System Prompt (如 "你是一位编程老师...")
   - 从 MySQL 加载历史消息 (滑动窗口 10 轮)
   - 组装: [系统提示词, 历史消息..., 用户新消息]
3. ModelService 获取模型客户端 (工厂模式 + 单例缓存)
4. DeepSeek/Qwen 流式返回 → JSON 编码每个 chunk
5. SSE 推送到浏览器 → marked.js 实时 Markdown 渲染
6. 完成后：消息持久化 MySQL, Token 统计, TTS 朗读
"""
c = add_card(s, 0.5, 2.0, 7.5, 5.0)
tf = card_text(c, flow_text.strip(), size=12, color=TEXT)

# Right side: key points
c2 = add_card(s, 8.5, 2.0, 4.3, 5.0)
tf2 = card_text(c2, "关键技术点", size=18, color=ACCENT, bold=True)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "JSON-SSE 换行保护", size=14, color=WHITE)
add_bullet(tf2, "chunk 用 json.dumps 编码\n换行符不再被 split('\\n') 吃掉\n保证 Markdown 语法完整", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "模型参数面板", size=14, color=WHITE)
add_bullet(tf2, "Temperature / MaxTokens / Top-P\n实时传入 ModelService\n覆盖配置文件默认值", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "停止生成按钮", size=14, color=WHITE)
add_bullet(tf2, "AbortController + fetch signal\n流式输出中途打断", size=11, color=TEXT_DIM)

# ====== SLIDE 7: TTS 语音 ======
s = dark_slide()
add_title(s, "语音功能 (TTS/STT)", size=32)
add_subtitle(s, "3 引擎 11 音色，浏览器语音输入 + 服务端语音合成", size=16)

c = add_card(s, 0.5, 2.0, 6.0, 5.0)
tf = card_text(c, "🎤 语音输入 (STT)", size=18, color=ACCENT, bold=True)
add_bullet(tf, "浏览器 SpeechRecognition API", size=14, color=WHITE)
add_bullet(tf, "continuous=false: 说完自动停止", size=12, color=TEXT_DIM)
add_bullet(tf, "onresult → 实时填入输入框 → 自动发送", size=12, color=TEXT_DIM)
add_bullet(tf, "", size=8)
add_bullet(tf, "🔊 语音输出 (TTS)", size=18, color=ACCENT, bold=True)
add_bullet(tf, "火山引擎: BV700/BV701/BV001/BV401 等 8 种", size=12, color=TEXT_DIM)
add_bullet(tf, "腾讯云: 智瑜/智美/智聆精品 3 种", size=12, color=TEXT_DIM)
add_bullet(tf, "失败自动兜底: 浏览器 speechSynthesis", size=12, color=TEXT_DIM)

c2 = add_card(s, 7.0, 2.0, 5.8, 5.0)
tf2 = card_text(c2, "TTS 技术细节", size=18, color=ACCENT2, bold=True)
add_bullet(tf2, "", size=6)
add_bullet(tf2, "火山引擎 TTS", size=14, color=WHITE)
add_bullet(tf2, "POST openspeech.bytedance.com/api/v1/tts", size=11, color=TEXT_DIM)
add_bullet(tf2, "Bearer Token 认证, JSON body", size=11, color=TEXT_DIM)
add_bullet(tf2, "返回 base64 MP3 → 前端 Blob 播放", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "腾讯云 TTS", size=14, color=WHITE)
add_bullet(tf2, "tencentcloud-sdk-python-tts SDK", size=11, color=TEXT_DIM)
add_bullet(tf2, "SecretId/Key 认证", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "发送前文字清洗", size=14, color=WHITE)
add_bullet(tf2, "去 Markdown 语法, 只保留纯朗读文字", size=11, color=TEXT_DIM)
add_bullet(tf2, "腾讯云限制 150 字 → 自动截断", size=11, color=TEXT_DIM)

# ====== SLIDE 8: 数据库设计 ======
s = dark_slide()
add_title(s, "数据库设计", size=32)
add_subtitle(s, "MySQL 8.0 utf8mb4, 4 张核心表", size=16)

db_text = """conversations (会话表)
  id VARCHAR(36) PK     title VARCHAR(200)     user_id VARCHAR(64) INDEX
  model VARCHAR(32)     role_id VARCHAR(36)     created_at / updated_at

messages (消息表)
  id INT PK AUTO_INC    conversation_id FK → conversations.id
  role VARCHAR(16)      content TEXT            tokens INT

roles (角色表)
  id VARCHAR(36) PK     name VARCHAR(100)       prompt TEXT
  icon VARCHAR(10)      is_builtin BOOLEAN      created_at / updated_at

users (用户表)
  id VARCHAR(64) PK     api_key VARCHAR(64) UNIQUE
  quota INT             created_at
"""
c = add_card(s, 0.5, 2.0, 7.5, 5.0)
tf = card_text(c, db_text, size=11, color=TEXT)

c2 = add_card(s, 8.5, 2.0, 4.3, 5.0)
tf2 = card_text(c2, "关键设计", size=18, color=ACCENT, bold=True)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "SQLAlchemy ORM", size=14, color=WHITE)
add_bullet(tf2, "Python 操作数据库\n无需手写 SQL", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "自动建库建表", size=14, color=WHITE)
add_bullet(tf2, "首次启动 create_all()\n自动创建所有表", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Token 估算", size=14, color=WHITE)
add_bullet(tf2, "中文字符 ×1.5\n英文单词 ×0.3\n数字 ×0.5", size=11, color=TEXT_DIM)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "硬删除机制", size=14, color=WHITE)
add_bullet(tf2, "先删 messages\n再删 conversation", size=11, color=TEXT_DIM)

# ====== SLIDE 9: 前端亮点 ======
s = dark_slide()
add_title(s, "前端技术亮点", size=32)
add_subtitle(s, "单文件 SPA (1900行) — 零框架依赖，纯原生 HTML/CSS/JS", size=16)

fe_items = [
    ("实时 Markdown", "marked.js + highlight.js 流式渲染\nJSON-SSE 换行保护，代码块语法高亮"),
    ("12 套主题", "CSS 变量驱动主题切换\n伪元素背景透明度 (文字不受影响)\n背景图/视频 + 历史记录"),
    ("粒子背景", "80 个 Canvas 粒子 + 近距离连线\nrequestAnimationFrame 60fps"),
    ("附件上传", "📎 点击上传 PDF/Word/TXT\nFormData → 后端解析 → 注入 AI 上下文"),
    ("对话导出", "从 DOM 直接抓取消息 → 生成 .md 文件\nBlob 下载，不调后端"),
    ("停止生成", "AbortController 中断 fetch\n发送按钮 ↔ 停止按钮动态切换"),
]
for i, (title, desc) in enumerate(fe_items):
    row, col = divmod(i, 3)
    c = add_card(s, 0.5+col*4.1, 2.0+row*2.5, 3.8, 2.0)
    tf = card_text(c, title, size=16, color=ACCENT, bold=True)
    add_bullet(tf, desc, size=11, color=TEXT_DIM)

# ====== SLIDE 10: API 接口一览 ======
s = dark_slide()
add_title(s, "API 接口一览", size=32)
add_subtitle(s, "20+ RESTful 接口, 5 个蓝图分组", size=16)

api_text = """对话 (chat.py)
  POST /api/v1/chat              普通对话 (同步)
  POST /api/v1/chat/stream       流式对话 (SSE, 打字机效果)

会话管理 (conversation.py)
  GET/POST  /api/v1/conversations        列表 / 创建
  GET/PUT/DELETE /api/v1/conversations/{id}  详情 / 改名 / 删除
  GET  /api/v1/conversations/{id}/messages   消息列表
  GET  /api/v1/conversations/{id}/export     导出 Markdown

角色管理 (role.py)
  GET/POST  /api/v1/roles          列表 / 创建 (同名检测)
  PUT/DELETE /api/v1/roles/{id}    更新 / 删除

语音合成 (voice.py)
  POST /api/v1/tts                 文本转语音 (火山/腾讯)
  GET  /api/v1/voices              音色列表

系统 (system.py)
  GET /api/v1/health               健康检查   GET /api/v1/models    模型列表
  GET /api/v1/stats                Token 统计  GET /api/v1/tools    工具列表
  GET /api/v1/local-media?path=   本地文件代理 (背景图/视频)

上传 (upload.py)
  POST /api/v1/upload              文件上传 (PDF/Word/TXT)"""
c = add_card(s, 0.5, 2.0, 12.3, 5.0)
tf = card_text(c, api_text, size=11, color=TEXT)

# ====== SLIDE 11: 部署方案 ======
s = dark_slide()
add_title(s, "部署方案", size=32)
add_subtitle(s, "三种部署方式，适配不同场景", size=16)

deploy_items = [
    ("开发模式", "python run.py → http://localhost:5000\nFlask 内置服务器 + SocketIO\n适合: 个人使用 / 开发调试"),
    ("生产模式", "Gunicorn + gevent + Nginx\n多 Worker 进程，SSE 友好\nproxy_buffering off (SSE 必须)\nSupervisor 进程守护\n适合: 对外服务 / 高并发"),
    ("Docker", "docker compose up -d\n一键部署，环境隔离\nMySQL + Redis 独立容器\n适合: 快速部署 / 迁移"),
]
for i, (title, desc) in enumerate(deploy_items):
    c = add_card(s, 0.5+i*4.1, 2.0, 3.8, 5.0)
    tf = card_text(c, title, size=18, color=ACCENT2, bold=True)
    for line in desc.split("\n"):
        add_bullet(tf, line, size=12, color=TEXT_DIM)

# ====== SLIDE 12: 项目总结 ======
s = dark_slide()
add_title(s, "项目总结", top=1.5, size=40)
add_subtitle(s, "从零构建的完整 AI 应用 —— 技术栈全面，功能完善，代码可读", top=2.2, size=18)

summary_items = [
    "✓ 后端 40+ Python 文件，5 个蓝图 20+ API 接口",
    "✓ 前端单文件 SPA (~1900行)，原生 JS 零框架依赖",
    "✓ MySQL 4 表设计 + SQLAlchemy ORM + Redis 限流",
    "✓ LangChain 策略模式集成 DeepSeek + 通义千问",
    "✓ SSE 流式输出 + JSON 编码解决换行断链",
    "✓ 火山/腾讯 双引擎 TTS + 浏览器 STT 语音输入",
    "✓ 12 套 CSS 变量主题 + 全透明度自定义",
    "✓ 模型参数面板 + Token 统计 + 对话导出 + 文件上传",
    "✓ Gunicorn + Nginx 生产部署 + Docker 容器化",
    "✓ 完整中文注释 + 零基础入门指南文档",
]
c = add_card(s, 1.0, 3.0, 11.3, 3.5)
tf = card_text(c, "", size=6)
for item in summary_items:
    add_bullet(tf, item, size=16, color=TEXT)

add_subtitle(s, "GitHub: github.com/RookieZJW/PythonAssisant", top=6.8, size=14)

# ====== SAVE ======
output_path = "E:/PythonAssistant/项目技术汇报.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
